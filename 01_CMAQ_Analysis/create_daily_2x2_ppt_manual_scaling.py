"""
CMAQ Daily 2x2 Fire Impact Maps - PowerPoint Generator
MANUAL SCALING VERSION

This script creates daily 2x2 comparison maps showing fire impacts on air pollutants.
You manually define the color scale levels for each pollutant.

Author: Generated for CMAQ Fire Analysis
Date: 2025-01-18
"""

import os
import numpy as np
import matplotlib.pyplot as plt
import pyrsig
import pycno
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches
import warnings
warnings.filterwarnings('ignore')

# ============= USER CONFIGURATION =============
# Pollutant to analyze
POLLUTANT = 'PM25_TOT'  # Options: 'O3', 'PM25_TOT', 'CO', 'BENZENE', 'TOLUENE', 'PHENOL'

# Unit conversion
CONVERT_TO_UGM3 = False  # True: convert to μg/m³, False: keep native units (ppb for gases)

# Paths
BASE_DIR = r'D:/Raw_Data/CMAQ_Model/'
OUTPUT_DIR = r'C:\Users\smtku\OA_Evolution_Wildfires\01_CMAQ_Analysis\figures\daily_2x2_maps'

# Daily aggregation method
DAILY_METHOD = 'mean'  # Options: 'mean' or 'max'

# Start date (first day of data)
START_DATE = datetime(2023, 6, 1)
# ==============================================


# ============= POLLUTANT CONFIGURATION =============
POLLUTANT_CONFIG = {
    'O3': {
        'var_name': 'O3',
        'display_name': 'O₃',
        'native_units': 'ppb',
        'colormap_base': 'viridis',
        'colormap_delta': 'RdBu_r',
        'levels_base': [0, 10, 20, 30, 40, 50, 60, 70, 80],
        'levels_delta': [-10, -5, -2, -1, 0, 1, 2, 5, 10, 20],
        'levels_percent': [-50, -20, -10, -5, 0, 5, 10, 20, 50, 100],
        'can_be_negative': True,
        'molecular_weight': 48.00,
    },
    'PM25_TOT': {
        'var_name': 'PM25_TOT',
        'display_name': 'PM₂.₅',
        'native_units': 'μg/m³',
        'colormap_base': 'viridis',
        'colormap_delta': 'YlOrRd',
        # Updated to match notebook (Pye2025_FirePMevolution.ipynb) Panel B0, B2, B3
        'levels_base': [0, 10, 20, 30, 40, 50, 60, 70, 80],           # Panel B0
        'levels_delta': [-0.5, 0, 0.5, 1, 2, 3, 4, 5, 10, 20],        # Panel B2
        'levels_percent': [-1, 0, 1, 2, 3, 4, 5, 10, 50, 100],        # Panel B3 (max 100%)
        'can_be_negative': False,
        'molecular_weight': None,  # Already in μg/m³
    },
    'CO': {
        'var_name': 'CO',
        'display_name': 'CO',
        'native_units': 'ppb',
        'colormap_base': 'viridis',
        'colormap_delta': 'YlOrRd',
        'levels_base': [0, 100, 200, 300, 400, 500, 750, 1000, 2000],
        'levels_delta': [0, 10, 20, 50, 100, 200, 500, 1000, 2000],
        'levels_percent': [0, 10, 20, 50, 100, 200, 500, 1000, 2000],
        'can_be_negative': False,
        'molecular_weight': 28.01,
    },
    'BENZENE': {
        'var_name': 'BENZENE',
        'display_name': 'Benzene',
        'native_units': 'ppb',
        'colormap_base': 'viridis',
        'colormap_delta': 'YlOrRd',
        'levels_base': [0, 0.05, 0.1, 0.2, 0.5, 1, 2, 5, 10],
        'levels_delta': [0, 0.01, 0.05, 0.1, 0.2, 0.5, 1, 2, 5],
        'levels_percent': [0, 10, 20, 50, 100, 200, 500, 1000, 2000],
        'can_be_negative': False,
        'molecular_weight': 78.11,
    },
    'TOLUENE': {
        'var_name': 'TOLUENE',
        'display_name': 'Toluene',
        'native_units': 'ppb',
        'colormap_base': 'viridis',
        'colormap_delta': 'YlOrRd',
        'levels_base': [0, 0.05, 0.1, 0.2, 0.5, 1, 2, 5, 10],
        'levels_delta': [0, 0.01, 0.05, 0.1, 0.2, 0.5, 1, 2, 5],
        'levels_percent': [0, 10, 20, 50, 100, 200, 500, 1000, 2000],
        'can_be_negative': False,
        'molecular_weight': 92.14,
    },
    'PHENOL': {
        'var_name': 'PHENOL',
        'display_name': 'Phenol',
        'native_units': 'ppb',
        'colormap_base': 'viridis',
        'colormap_delta': 'YlOrRd',
        'levels_base': [0, 0.01, 0.05, 0.1, 0.2, 0.5, 1, 2, 5],
        'levels_delta': [0, 0.005, 0.01, 0.05, 0.1, 0.2, 0.5, 1, 2],
        'levels_percent': [0, 10, 20, 50, 100, 200, 500, 1000, 2000],
        'can_be_negative': False,
        'molecular_weight': 94.11,
    },
}
# ===================================================


def load_cmaq_data():
    """Load CMAQ base and no-fire data"""
    print("Loading CMAQ data...")

    cmaq_base_path = BASE_DIR + 'netcdffiles/COMBINE_ACONC_cmaq6acracmm3_base_2023_12US4_202306.nc'
    cmaq_nofire_path = BASE_DIR + 'netcdffiles/COMBINE_ACONC_cmaq6acracmm3_nofire_2023_12US4_202306.nc'
    metcro2d_path = BASE_DIR + 'MOD3DATA_MET/METCRO2D.12US4.35L.230701'
    gridcro2d_path = BASE_DIR + 'MOD3DATA_MET/GRIDCRO2D.12US4.35L.230701'

    baseconc = pyrsig.open_ioapi(cmaq_base_path)
    nofireconc = pyrsig.open_ioapi(cmaq_nofire_path)
    metcro2d = pyrsig.open_ioapi(metcro2d_path)

    # Load projection info
    cno = pycno.cno(baseconc.crs_proj4)

    print("Data loaded successfully!")
    return baseconc, nofireconc, metcro2d, cno


def calculate_fire_impact(baseconc, nofireconc, metcro2d, config):
    """Calculate fire impact (base - nofire) for the pollutant"""
    print(f"Calculating fire impact for {config['display_name']}...")

    var_name = config['var_name']

    # Get base and nofire concentrations
    base = baseconc[var_name]
    nofire = nofireconc[var_name]

    # Select surface layer (LAY=0) if data has layer dimension
    if len(base.shape) == 4:  # (TSTEP, LAY, ROW, COL)
        print("  Selecting surface layer (LAY=0)...")
        base = base[:, 0, :, :]
        nofire = nofire[:, 0, :, :]

    # Convert to μg/m³ if requested and if it's a gas (ppb)
    if CONVERT_TO_UGM3 and config['native_units'] == 'ppb':
        print(f"Converting from {config['native_units']} to μg/m³...")
        # AIR_DENS is in the COMBINE_ACONC file, not METCRO2D
        air_dens = baseconc['AIR_DENS']  # Air density (kg/m³)

        # Select surface layer for air density if needed
        if len(air_dens.shape) == 4:
            air_dens = air_dens[:, 0, :, :]

        mw = config['molecular_weight']

        # Convert: μg/m³ = ppb × air_dens × MW / 28.9628
        base = base * air_dens * mw / 28.9628
        nofire = nofire * air_dens * mw / 28.9628

        units = 'μg/m³'
    else:
        units = config['native_units']

    # Calculate delta
    delta = base - nofire

    print(f"Fire impact calculated! Units: {units}")
    return base, nofire, delta, units


def aggregate_to_daily(data, method='mean'):
    """Aggregate hourly data to daily"""
    print(f"Aggregating to daily {method}...")

    # Reshape from (720, ROW, COL) to (30, 24, ROW, COL)
    n_days = 30
    n_hours = 24

    # Get data array
    data_array = data.values

    # Reshape
    shape = (n_days, n_hours) + data_array.shape[1:]
    data_reshaped = data_array.reshape(shape)

    # Aggregate
    if method == 'mean':
        daily_data = np.mean(data_reshaped, axis=1)
    elif method == 'max':
        daily_data = np.max(data_reshaped, axis=1)
    else:
        raise ValueError(f"Unknown method: {method}")

    print(f"Daily aggregation complete! Shape: {daily_data.shape}")
    return daily_data


def format_map_axis(ax, cno, title):
    """Format map axis"""
    ax.set_xlabel('')
    ax.set_xticks([])
    ax.set_ylabel('')
    ax.set_yticks([])
    ax.set_title(title, fontsize=10)
    cno.drawstates(ax=ax, linewidth=0.2)


def create_daily_plot(base_day, nofire_day, delta_day, config, units, date_str, cno, day_num):
    """Create 2x2 plot for a single day"""

    fig, axes = plt.subplots(2, 2, figsize=(14, 7), dpi=150)
    plt.subplots_adjust(wspace=0.25, hspace=0.35, top=0.92, bottom=0.08, left=0.05, right=0.95)

    # Overall title - simplified to just pollutant name and date
    fig.suptitle(f'{config["display_name"]} - {date_str}',
                 fontsize=16, fontweight='bold')

    # Panel (a): Base simulation
    pv = axes[0, 0].contourf(base_day, cmap=config['colormap_base'],
                              levels=config['levels_base'], extend='both')
    cb = plt.colorbar(pv, ax=axes[0, 0])
    cb.set_label(f'{config["display_name"]} ({units})', rotation=270, labelpad=15, fontsize=9)
    title_a = f'(a) Base Simulation\nMean: {np.nanmean(base_day):.2f}, Max: {np.nanmax(base_day):.1f} {units}'
    format_map_axis(axes[0, 0], cno, title_a)

    # Panel (b): No-fire simulation
    pv = axes[0, 1].contourf(nofire_day, cmap=config['colormap_base'],
                              levels=config['levels_base'], extend='both')
    cb = plt.colorbar(pv, ax=axes[0, 1])
    cb.set_label(f'{config["display_name"]} ({units})', rotation=270, labelpad=15, fontsize=9)
    title_b = f'(b) No-Fire Scenario\nMean: {np.nanmean(nofire_day):.2f}, Max: {np.nanmax(nofire_day):.1f} {units}'
    format_map_axis(axes[0, 1], cno, title_b)

    # Panel (c): Absolute delta
    pv = axes[1, 0].contourf(delta_day, cmap=config['colormap_delta'],
                              levels=config['levels_delta'], extend='both')
    cb = plt.colorbar(pv, ax=axes[1, 0])
    cb.set_label(f'Δ{config["display_name"]} ({units})', rotation=270, labelpad=15, fontsize=9)
    title_c = f'(c) Fire Impact: Δ{config["display_name"]} (Base − No Fire)\nMean: {np.nanmean(delta_day):.2f}, Range: [{np.nanmin(delta_day):.2f}, {np.nanmax(delta_day):.2f}] {units}'
    format_map_axis(axes[1, 0], cno, title_c)

    # Panel (d): Percent delta
    # Fix: Use base_day as denominator (not nofire_day) and use pre-calculated delta_day
    # This matches the reference notebook methodology and handles max aggregation correctly
    percent_delta = np.where(base_day > 0, (delta_day / base_day * 100), np.nan)

    pv = axes[1, 1].contourf(percent_delta, cmap=config['colormap_delta'],
                              levels=config['levels_percent'], extend='both')
    cb = plt.colorbar(pv, ax=axes[1, 1])
    cb.set_label('% Change', rotation=270, labelpad=15, fontsize=9)
    title_d = f'(d) Relative Fire Impact: %Δ{config["display_name"]}\nMean: {np.nanmean(percent_delta):.1f}%, Range: [{np.nanmin(percent_delta):.0f}%, {np.nanmax(percent_delta):.0f}%]'
    format_map_axis(axes[1, 1], cno, title_d)

    # Save to temporary file
    temp_file = f'temp_day_{day_num:02d}.png'
    plt.savefig(temp_file, dpi=150, bbox_inches='tight')
    plt.close()

    return temp_file


def create_powerpoint(image_files, dates, output_path):
    """Create PowerPoint presentation with daily maps"""
    print("Creating PowerPoint presentation...")

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # Widescreen (16:9)
    prs.slide_height = Inches(7.5)

    for img_file, date_str in zip(image_files, dates):
        # Add blank slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add image (full slide)
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(12.333)
        slide.shapes.add_picture(img_file, left, top, width=width)

    # Save presentation
    prs.save(output_path)
    print(f"PowerPoint saved: {output_path}")

    # Clean up temporary files
    for img_file in image_files:
        if os.path.exists(img_file):
            os.remove(img_file)
    print("Temporary files cleaned up.")


def main():
    """Main execution function"""
    print("="*60)
    print("CMAQ Daily 2x2 Fire Impact Maps Generator")
    print("Manual Scaling Version")
    print("="*60)

    # Validate pollutant
    if POLLUTANT not in POLLUTANT_CONFIG:
        raise ValueError(f"Unknown pollutant: {POLLUTANT}. Options: {list(POLLUTANT_CONFIG.keys())}")

    config = POLLUTANT_CONFIG[POLLUTANT]
    print(f"\nPollutant: {config['display_name']}")
    print(f"Native units: {config['native_units']}")
    print(f"Convert to μg/m³: {CONVERT_TO_UGM3}")
    print(f"Daily aggregation: {DAILY_METHOD}")

    # Create output directory
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    # Load data
    baseconc, nofireconc, metcro2d, cno = load_cmaq_data()

    # Calculate fire impact
    base, nofire, delta, units = calculate_fire_impact(baseconc, nofireconc, metcro2d, config)

    # Aggregate to daily
    base_daily = aggregate_to_daily(base, method=DAILY_METHOD)
    nofire_daily = aggregate_to_daily(nofire, method=DAILY_METHOD)

    # IMPORTANT: For 'max' aggregation, delta must be calculated AFTER aggregation
    # because max(base - nofire) ≠ max(base) - max(nofire)
    # For 'mean' aggregation, it doesn't matter: mean(base - nofire) = mean(base) - mean(nofire)
    if DAILY_METHOD == 'max':
        # Calculate delta from aggregated values (correct for max)
        delta_daily = base_daily - nofire_daily
        print("  Delta calculated from aggregated base and nofire (correct for max method)")
    else:
        # For mean, we can aggregate delta directly (mathematically equivalent)
        delta_daily = aggregate_to_daily(delta, method=DAILY_METHOD)

    # Generate plots for each day
    print("\nGenerating daily plots...")
    image_files = []
    dates = []

    for day in range(30):
        date = START_DATE + timedelta(days=day)
        date_str = date.strftime('%Y-%m-%d')
        dates.append(date_str)

        print(f"  Processing {date_str} (Day {day+1}/30)...")

        img_file = create_daily_plot(
            base_daily[day],
            nofire_daily[day],
            delta_daily[day],
            config,
            units,
            date_str,
            cno,
            day + 1
        )
        image_files.append(img_file)

    # Create output filename
    unit_str = units.replace('/', 'per').replace('³', '3').replace('μg', 'ug')
    output_filename = f"{config['var_name']}_{unit_str}_daily_2x2_June2023.pptx"
    output_path = os.path.join(OUTPUT_DIR, output_filename)

    # Create PowerPoint
    create_powerpoint(image_files, dates, output_path)

    print("\n" + "="*60)
    print("COMPLETE!")
    print(f"Output saved to: {output_path}")
    print("="*60)


if __name__ == "__main__":
    main()
